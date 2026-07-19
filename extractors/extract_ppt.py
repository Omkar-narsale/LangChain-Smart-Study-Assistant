"""
extract_ppt.py — PPT/PPTX Text Extractor
========================================
Extracts text from PowerPoint presentations (slide titles, bullet points, notes, tables).
"""

from __future__ import annotations
import io
import logging

from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_ppt(file: io.IOBase) -> tuple[str, int]:
    """
    Extract text from a PPT/PPTX document.
    Returns a tuple of (extracted_text, slide_count).
    """
    try:
        prs = Presentation(file)
        slide_count = len(prs.slides)
        slides_text = []

        for i, slide in enumerate(prs.slides):
            slide_content = []
            
            # Extract notes if available
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_content.append(f"Speaker Notes:\n{notes}")

            # Extract shapes (text boxes, tables)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_content.append(shape.text.strip())
                    
                if shape.has_table:
                    table_content = []
                    for row in shape.table.rows:
                        row_data = [cell.text.strip() for cell in row.cells if cell.text]
                        if row_data:
                            table_content.append(" | ".join(row_data))
                    if table_content:
                        slide_content.append("Table Data:\n" + "\n".join(table_content))
            
            if slide_content:
                slides_text.append(f"--- Slide {i+1} ---\n" + "\n\n".join(slide_content))

        return "\n\n".join(slides_text), slide_count

    except Exception as exc:
        logger.exception("PPT extraction failed.")
        raise RuntimeError(f"PPT text extraction failed: {exc}") from exc
