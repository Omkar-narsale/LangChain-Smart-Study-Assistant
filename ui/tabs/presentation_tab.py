"""
ui/tabs/presentation_tab.py — AI Presentation Studio Tab
=========================================================
Independent flagship module allowing users to convert documents/text/chat history
into premium, presentation-ready PowerPoint and PDF slides.
"""

from __future__ import annotations
import json
import re
import io
import time
import streamlit as st
import streamlit.components.v1 as components
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from config import get_model, task_token_budget, invoke_with_retry
from ui.components import render_empty_state

# ---------------------------------------------------------------------------
# PDF Slide Builder (using ReportLab)
# ---------------------------------------------------------------------------

def generate_pdf_file(outline: list, settings: dict) -> bytes:
    """Generate slide PDF representation using reportlab."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.colors import HexColor

    colors_map = {
        "Blue": ("#0B1220", "#00C2FF", "#F1F5F9", "#94A3B8"),
        "Purple": ("#080E1A", "#7C5CFF", "#F1F5F9", "#94A3B8"),
        "Green": ("#060F0E", "#22C55E", "#F1F5F9", "#94A3B8"),
        "Orange": ("#0F0A06", "#F97316", "#F1F5F9", "#94A3B8"),
        "Red": ("#0F0606", "#EF4444", "#F1F5F9", "#94A3B8"),
        "Auto": ("#0B1220", "#00C2FF", "#F1F5F9", "#94A3B8")
    }

    palette = settings.get("color_palette", "Auto")
    bg_hex, prim_hex, text_hex, sec_hex = colors_map.get(palette, colors_map["Auto"])

    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=(800, 450)) # 16:9 Landscape

    total_slides = len(outline)

    for i, slide in enumerate(outline):
        layout = slide.get("layout", "Content")
        title = slide.get("title", "")
        bullets = slide.get("bullets", [])

        # Background color fill
        pdf.setFillColor(HexColor(bg_hex))
        pdf.rect(0, 0, 800, 450, fill=True, stroke=False)

        if layout in ["Title", "Section Divider"]:
            # Decorative left panel
            pdf.setFillColor(HexColor(prim_hex))
            pdf.rect(0, 0, 16, 450, fill=True, stroke=False)

            # Centered Text
            pdf.setFillColor(HexColor(text_hex))
            pdf.setFont("Helvetica-Bold", 36)
            pdf.drawCentredString(400, 240, title)

            if layout == "Title":
                pdf.setFillColor(HexColor(prim_hex))
                pdf.setFont("Helvetica", 16)
                pdf.drawCentredString(400, 180, settings.get("type", "Presentation"))
        else:
            # Header title
            pdf.setFillColor(HexColor(prim_hex))
            pdf.setFont("Helvetica-Bold", 26)
            pdf.drawString(40, 390, title)

            # Divider Line
            pdf.setStrokeColor(HexColor(prim_hex))
            pdf.setLineWidth(2)
            pdf.line(40, 375, 760, 375)

            # Draw bullets or columns
            pdf.setFillColor(HexColor(text_hex))
            pdf.setFont("Helvetica", 14)
            y = 320
            
            if layout == "Comparison" and len(bullets) >= 2:
                # 2-column rendering
                mid = len(bullets) // 2
                col1 = bullets[:mid]
                col2 = bullets[mid:]
                
                # Column 1
                y1 = 320
                for bullet in col1:
                    pdf.drawString(50, y1, f"•  {bullet}")
                    y1 -= 35
                
                # Column 2
                y2 = 320
                for bullet in col2:
                    pdf.drawString(420, y2, f"•  {bullet}")
                    y2 -= 35
            else:
                # Standard bullet rendering
                for bullet in bullets:
                    pdf.drawString(50, y, f"•  {bullet}")
                    y -= 35
                    if y < 70:
                        break

            # Footer
            pdf.setFillColor(HexColor(sec_hex))
            pdf.setFont("Helvetica", 10)
            pdf.drawString(40, 25, f"⚡ Study AI | {settings.get('type', 'Presentation')}")
            pdf.drawRightString(760, 25, f"{i+1} / {total_slides}")

        pdf.showPage()

    pdf.save()
    buffer.seek(0)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# PPT Slide Builder (using python-pptx)
# ---------------------------------------------------------------------------

def generate_pptx_file(outline: list, settings: dict) -> bytes:
    """Generate pptx presentation file binary."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    colors_map = {
        "Blue": ("0B1220", "00C2FF", "F1F5F9", "94A3B8"),
        "Purple": ("080E1A", "7C5CFF", "F1F5F9", "94A3B8"),
        "Green": ("060F0E", "22C55E", "F1F5F9", "94A3B8"),
        "Orange": ("0F0A06", "F97316", "F1F5F9", "94A3B8"),
        "Red": ("0F0606", "EF4444", "F1F5F9", "94A3B8"),
        "Auto": ("0B1220", "00C2FF", "F1F5F9", "94A3B8")
    }

    palette = settings.get("color_palette", "Auto")
    bg_hex, prim_hex, text_hex, sec_hex = colors_map.get(palette, colors_map["Auto"])

    def hex_to_rgb(hex_str):
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    bg_rgb = hex_to_rgb(bg_hex)
    prim_rgb = hex_to_rgb(prim_hex)
    text_rgb = hex_to_rgb(text_hex)
    sec_rgb = hex_to_rgb(sec_hex)

    total_slides = len(outline)

    for i, slide_data in enumerate(outline):
        slide = prs.slides.add_slide(blank_layout)

        # Full background solid color fill
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()

        layout = slide_data.get("layout", "Content")
        title_text = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes", "")

        # Attach presenter notes if selected
        if settings.get("include_notes", True) and notes:
            slide.notes_slide.notes_text_frame.text = notes

        if layout in ["Title", "Section Divider"]:
            # Elegant vertical accent line on the left side
            accent = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = prim_rgb
            accent.line.fill.background()

            # Creative block on the right half (asymmetrical design)
            right_block = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(9.5), Inches(0), Inches(3.833), Inches(7.5)
            )
            right_block.fill.solid()
            right_block.fill.fore_color.rgb = prim_rgb
            right_block.line.fill.background()

            # Main Slide Text Area
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(3.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.LEFT
            p.font.name = "Arial"
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = text_rgb

            if layout == "Title":
                p2 = tf.add_paragraph()
                p2.text = settings.get("type", "Presentation")
                p2.alignment = PP_ALIGN.LEFT
                p2.font.name = "Arial"
                p2.font.size = Pt(20)
                p2.font.color.rgb = prim_rgb
                p2.space_before = Pt(20)
        else:
            # Header Row
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.73), Inches(1.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = "Arial"
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = prim_rgb

            # Creative top-right accent dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(12.2), Inches(0.6), Inches(0.25), Inches(0.25)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = prim_rgb
            dot.line.fill.background()

            # Header Accent Divider
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.73), Inches(0.03)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = prim_rgb
            line.line.fill.background()

            # Footer Logo
            footerBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(5.0), Inches(0.4))
            ftf = footerBox.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"⚡ Study AI | {settings.get('type', 'Presentation')}"
            fp.font.name = "Arial"
            fp.font.size = Pt(10)
            fp.font.color.rgb = sec_rgb

            # Slide Indicator Count
            numBox = slide.shapes.add_textbox(Inches(11.5), Inches(6.8), Inches(1.0), Inches(0.4))
            ntf = numBox.text_frame
            np = ntf.paragraphs[0]
            np.text = f"{i+1} / {total_slides}"
            np.alignment = PP_ALIGN.RIGHT
            np.font.name = "Arial"
            np.font.size = Pt(10)
            np.font.color.rgb = sec_rgb

            # Standard Content Slide Layouts
            if layout == "Comparison" and len(bullets) >= 2:
                mid = len(bullets) // 2
                
                # Left pane content
                col1_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.5))
                tf1 = col1_box.text_frame
                tf1.word_wrap = True
                for idx, b in enumerate(bullets[:mid]):
                    p_b = tf1.add_paragraph() if idx > 0 else tf1.paragraphs[0]
                    p_b.text = f"•  {b}"
                    p_b.font.name = "Arial"
                    p_b.font.size = Pt(16)
                    p_b.font.color.rgb = text_rgb
                    p_b.space_after = Pt(10)

                # Right pane content
                col2_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(5.6), Inches(4.5))
                tf2 = col2_box.text_frame
                tf2.word_wrap = True
                for idx, b in enumerate(bullets[mid:]):
                    p_b = tf2.add_paragraph() if idx > 0 else tf2.paragraphs[0]
                    p_b.text = f"•  {b}"
                    p_b.font.name = "Arial"
                    p_b.font.size = Pt(16)
                    p_b.font.color.rgb = text_rgb
                    p_b.space_after = Pt(10)

            elif layout in ["Timeline", "Process"]:
                # Draw elegant connected pipeline block design
                steps = min(len(bullets), 4)
                if steps == 0:
                    steps = 3
                    bullets = ["Step 1", "Step 2", "Step 3"]
                
                w_per_block = 11.73 / steps
                for idx, txt in enumerate(bullets[:steps]):
                    l_pos = 0.8 + (idx * w_per_block) + 0.1
                    w_block = w_per_block - 0.2

                    card = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l_pos), Inches(2.2), Inches(w_block), Inches(3.8)
                    )
                    card.fill.solid()
                    card.fill.fore_color.rgb = hex_to_rgb(bg_hex)
                    card.line.color.rgb = prim_rgb
                    card.line.width = Pt(1.5)

                    ctf = card.text_frame
                    ctf.word_wrap = True
                    cp = ctf.paragraphs[0]
                    cp.text = f"Phase {idx+1}"
                    cp.font.name = "Arial"
                    cp.font.size = Pt(16)
                    cp.font.bold = True
                    cp.font.color.rgb = prim_rgb
                    cp.alignment = PP_ALIGN.CENTER
                    cp.space_after = Pt(10)

                    cp2 = ctf.add_paragraph()
                    cp2.text = txt
                    cp2.font.name = "Arial"
                    cp2.font.size = Pt(13)
                    cp2.font.color.rgb = text_rgb
                    cp2.alignment = PP_ALIGN.CENTER
            else:
                # Default Standard Bullet list slide content
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.73), Inches(4.5))
                ctf = content_box.text_frame
                ctf.word_wrap = True
                for idx, b in enumerate(bullets):
                    p_b = ctf.add_paragraph() if idx > 0 else ctf.paragraphs[0]
                    p_b.text = f"•  {b}"
                    p_b.font.name = "Arial"
                    p_b.font.size = Pt(18)
                    p_b.font.color.rgb = text_rgb
                    p_b.space_after = Pt(12)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


# ---------------------------------------------------------------------------
# Outline Parser
# ---------------------------------------------------------------------------

def _parse_outline_json(raw: str) -> list:
    """Parse output from LLM safely into a structured slide list."""
    cleaned = raw.strip()
    cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
    cleaned = re.sub(r"\s*```$", "", cleaned)
    cleaned = cleaned.strip()

    try:
        return json.loads(cleaned)
    except Exception:
        pass

    match = re.search(r"\[\s*\{.*\}\s*\]", cleaned, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except Exception:
            pass

    # Standard fallback outline
    return [
        {"title": "Title Slide", "bullets": ["Overview of this presentation"], "notes": "Welcome to the presentation.", "layout": "Title"},
        {"title": "Introduction", "bullets": ["Overview of key topics", "Core concepts list"], "notes": "Let's introduce the theme.", "layout": "Content"},
        {"title": "Summary", "bullets": ["Key takeaways", "Next steps to implement"], "notes": "Wrapping up the main findings.", "layout": "Conclusion"}
    ]


# ---------------------------------------------------------------------------
# LLM Outline Generation Call
# ---------------------------------------------------------------------------

def generate_outline_llm(source_text: str, settings: dict) -> list:
    """Call active LLM to generate outline JSON list."""
    from langchain_core.output_parsers import StrOutputParser
    from langchain_core.prompts import PromptTemplate

    prompt = PromptTemplate.from_template(
        "You are a professional presentation planner. Output ONLY a valid JSON array of slide objects in {language}.\n"
        "Do NOT write any preamble, introduction, markdown blocks (such as ```json), or explanation. Output raw JSON only.\n\n"
        "SETTINGS:\n"
        "- Presentation Type: {type}\n"
        "- Total Slides Needed: {slides}\n"
        "- Layouts allowed: 'Title', 'Agenda', 'Section Divider', 'Content', 'Comparison', 'Timeline', 'Process', 'Conclusion', 'Thank You'\n"
        "- Include Speaker Notes: {include_notes}\n\n"
        "OUTLINE CONSTRAINTS:\n"
        "- Maximum 5-7 bullet points per slide.\n"
        "- Maximum 15 words per bullet point.\n"
        "- Presenter notes should be 1-3 professional sentences detailing slide contents.\n\n"
        "SOURCE CONTEXT:\n"
        "{context}\n\n"
        "JSON output format example:\n"
        "[\n"
        "  {{\n"
        "    \"title\": \"Slide Title\",\n"
        "    \"bullets\": [\"Bullet point 1\", \"Bullet point 2\"],\n"
        "    \"notes\": \"Speaker note here.\",\n"
        "    \"layout\": \"Content\"\n"
        "  }}\n"
        "]"
    )

    parser = StrOutputParser()
    with task_token_budget("study_material") as model:
        chain = prompt | model | parser
        raw_json = invoke_with_retry(
            chain,
            {
                "language": settings.get("language", "English"),
                "type": settings.get("type", "Business"),
                "slides": settings.get("slides", 6),
                "include_notes": "Yes" if settings.get("include_notes", True) else "No",
                "context": source_text[:12000] # Fit model token limits
            }
        )

    return _parse_outline_json(raw_json)


# ---------------------------------------------------------------------------
# Main Render Function
# ---------------------------------------------------------------------------

def render_presentation_tab() -> None:
    st.markdown("<h2 class='text-section'>🎨 AI Presentation Studio</h2>", unsafe_allow_html=True)
    st.markdown("<p class='text-body'>Design, customize, and build premium presentations from any source of knowledge.</p>", unsafe_allow_html=True)

    # ── Source Selection Detection ──────────────────────────────────────────
    file_name = st.session_state.get("file_display_name", "")
    has_pdf = file_name.lower().endswith(".pdf")
    has_docx = file_name.lower().endswith((".docx", ".doc"))
    has_ppt = file_name.lower().endswith((".pptx", ".ppt"))
    has_study_notes = st.session_state.get("study_output") is not None
    has_chat = len(st.session_state.get("chat_history", [])) > 0

    st.markdown("<h3 class='text-card-title' style='margin-top:20px;'>1. Choose Input Source</h3>", unsafe_allow_html=True)

    source_cols = st.columns(6)
    sources = [
        {"id": "Paste Text", "icon": "📋", "label": "Paste Text", "available": True},
        {"id": "PDF", "icon": "📄", "label": "PDF", "available": has_pdf},
        {"id": "DOCX", "icon": "📝", "label": "DOCX", "available": has_docx},
        {"id": "Existing PPT", "icon": "📊", "label": "Existing PPT", "available": has_ppt},
        {"id": "Study Notes", "icon": "📚", "label": "Study Notes", "available": has_study_notes},
        {"id": "AI Chat", "icon": "💬", "label": "AI Chat", "available": has_chat},
    ]

    selected_source = st.session_state.get("presentation_source", "Paste Text")

    for idx, src in enumerate(sources):
        with source_cols[idx]:
            if not src["available"]:
                # Disabled visual state card
                st.markdown(f"""
                <div style="
                    border: 1px solid rgba(255,255,255,0.03);
                    background: rgba(255,255,255,0.01);
                    border-radius: 12px;
                    padding: 16px 8px;
                    text-align: center;
                    opacity: 0.35;
                    cursor: not-allowed;
                ">
                    <div style="font-size:20px;">{src['icon']}</div>
                    <div style="font-size:11px;font-weight:600;margin-top:6px;color:var(--text-muted);">{src['label']}</div>
                </div>
                """, unsafe_allow_html=True)
            else:
                # Selectable state
                is_active = (selected_source == src["id"])
                border_style = "1px solid var(--primary)" if is_active else "1px solid var(--border)"
                bg_style = "linear-gradient(135deg, rgba(124,92,255,0.12), rgba(124,92,255,0.03))" if is_active else "var(--glass-bg)"
                shadow_style = "var(--shadow-glow)" if is_active else "none"

                st.markdown(f"""
                <div style="
                    border: {border_style};
                    background: {bg_style};
                    border-radius: 12px;
                    padding: 16px 8px;
                    text-align: center;
                    box-shadow: {shadow_style};
                    cursor: pointer;
                ">
                    <div style="font-size:20px;">{src['icon']}</div>
                </div>
                """, unsafe_allow_html=True)
                
                # Active button underneath the custom card
                if st.button(src["label"], key=f"src_btn_{src['id']}", use_container_width=True):
                    st.session_state.presentation_source = src["id"]
                    st.rerun()

    # Paste text area rendering
    if st.session_state.presentation_source == "Paste Text":
        pasted = st.text_area(
            "Paste presentation source text here:",
            value=st.session_state.get("presentation_pasted_text", ""),
            height=200,
            key="presentation_pasted_text_input"
        )
        st.session_state.presentation_pasted_text = pasted
    else:
        st.info(f"Source context selected: **{st.session_state.presentation_source}**")

    # ── Presentation Settings (Section 2) ──────────────────────────────────
    st.markdown("<h3 class='text-card-title' style='margin-top:32px;'>2. Presentation Settings</h3>", unsafe_allow_html=True)

    with st.container():
        st.markdown('<div class="glass-card" style="padding:24px; margin-bottom:20px;">', unsafe_allow_html=True)
        set_c1, set_c2 = st.columns(2)

        with set_c1:
            pres_type = st.selectbox(
                "Presentation Type",
                ["Academic", "Business", "Seminar", "Research", "Project Presentation", "Pitch Deck", "Workshop", "Training"],
                index=1
            )
            slide_count = st.slider("Slides Count", min_value=5, max_value=20, value=8)

            theme = st.selectbox(
                "Design Theme",
                ["Modern", "Glass", "Minimal", "Corporate", "Gradient", "Dark Professional"],
                index=5
            )

        with set_c2:
            palette = st.selectbox(
                "Color Palette",
                ["Blue", "Purple", "Green", "Orange", "Red", "Auto"],
                index=0
            )

            image_style = st.selectbox(
                "Visual Elements Style",
                ["Icons only", "Illustrations", "AI diagrams", "Minimal", "No Images"],
                index=0
            )

            language = st.selectbox(
                "Presentation Language",
                ["English", "Hindi", "Marathi"],
                index=0
            )

        include_notes = st.toggle("Include Speaker Notes", value=True)
        st.markdown('</div>', unsafe_allow_html=True)

    # ── Generate Outline Button ───────────────────────────────────────────
    if st.button("📝 Generate Slide Outline", type="primary", use_container_width=True):
        # Extract content text based on source selection
        src = st.session_state.presentation_source
        source_text = ""

        if src == "Paste Text":
            source_text = st.session_state.get("presentation_pasted_text", "").strip()
        elif src in ["PDF", "DOCX", "Existing PPT"]:
            source_text = st.session_state.get("document_text", "").strip()
        elif src == "Study Notes" and has_study_notes:
            out = st.session_state.study_output
            summary_content = out.get("summary", "")
            kp_content = "\n".join(out.get("key_points", []))
            source_text = f"SUMMARY:\n{summary_content}\n\nKEY POINTS:\n{kp_content}"
        elif src == "AI Chat" and has_chat:
            chat_lines = []
            for m in st.session_state.chat_history:
                role = "User" if m["role"] == "user" else "AI"
                chat_lines.append(f"{role}: {m['content']}")
            source_text = "\n".join(chat_lines)

        if not source_text:
            st.error("⚠️ Source text is empty. Please select or paste text first.")
            return

        with st.spinner("Generating Outline..."):
            settings = {
                "type": pres_type,
                "slides": slide_count,
                "theme": theme,
                "color_palette": palette,
                "image_style": image_style,
                "language": language,
                "include_notes": include_notes
            }
            try:
                outline = generate_outline_llm(source_text, settings)
                st.session_state.presentation_outline = outline
                st.session_state.presentation_settings = settings
                st.session_state.presentation_file = None # Clear old slides
            except Exception as e:
                st.error(f"Outline generation failed: {e}")

    # ── Preview Outline & Edit (Section 3) ──────────────────────────────────
    if st.session_state.get("presentation_outline"):
        st.markdown("<h3 class='text-card-title' style='margin-top:32px;'>3. Preview & Edit Slides Content</h3>", unsafe_allow_html=True)
        st.markdown("<p class='text-body'>You can edit slide titles, bullet points, and speaker notes directly before compiling the presentation file.</p>", unsafe_allow_html=True)

        outline = st.session_state.presentation_outline

        for idx, slide in enumerate(outline):
            st.markdown(f'<div class="glass-card" style="padding:18px; margin-bottom:12px;">', unsafe_allow_html=True)
            st.markdown(f"**Slide {idx+1} ({slide.get('layout', 'Content')} Layout)**")

            # Edit slide Title
            new_title = st.text_input(
                f"Slide {idx+1} Title",
                value=slide.get("title", ""),
                key=f"slide_title_inp_{idx}"
            )
            slide["title"] = new_title

            # Edit slide Bullet points
            bullets_raw = "\n".join(slide.get("bullets", []))
            new_bullets = st.text_area(
                f"Slide {idx+1} Bullets (one per line)",
                value=bullets_raw,
                height=110,
                key=f"slide_bullets_inp_{idx}"
            )
            slide["bullets"] = [b.strip() for b in new_bullets.split("\n") if b.strip()]

            # Edit speaker notes
            if slide.get("notes") is not None:
                new_notes = st.text_area(
                    f"Slide {idx+1} Presenter Notes",
                    value=slide.get("notes", ""),
                    height=70,
                    key=f"slide_notes_inp_{idx}"
                )
                slide["notes"] = new_notes

            st.markdown('</div>', unsafe_allow_html=True)

        st.session_state.presentation_outline = outline

        # ── Export & Download (Section 4) ───────────────────────────────────
        st.markdown("<h3 class='text-card-title' style='margin-top:32px;'>4. Export Presentation</h3>", unsafe_allow_html=True)

        exp_c1, exp_c2 = st.columns(2)

        # Content signature definition for caching
        outline_sig = json.dumps(st.session_state.presentation_outline) + json.dumps(st.session_state.presentation_settings)

        with exp_c1:
            if st.button("🎨 Create Presentation Files", type="primary", use_container_width=True):
                with st.spinner("Designing Slides & Compiling PPTX..."):
                    # Check cache first
                    cache = st.session_state.get("presentation_cache", {})
                    if outline_sig in cache:
                        st.session_state.presentation_file = cache[outline_sig]
                    else:
                        pptx_bin = generate_pptx_file(
                            st.session_state.presentation_outline,
                            st.session_state.presentation_settings
                        )
                        pdf_bin = generate_pdf_file(
                            st.session_state.presentation_outline,
                            st.session_state.presentation_settings
                        )
                        st.session_state.presentation_file = {
                            "pptx": pptx_bin,
                            "pdf": pdf_bin
                        }
                        st.session_state.presentation_cache[outline_sig] = st.session_state.presentation_file

                    st.success("✅ Presentation designed successfully!")

        if st.session_state.get("presentation_file"):
            files = st.session_state.presentation_file
            
            with exp_c2:
                # PPTX Download
                st.download_button(
                    "⬇ Download PPTX",
                    data=files["pptx"],
                    file_name="presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                # PDF Download
                st.download_button(
                    "⬇ Download PDF Version",
                    data=files["pdf"],
                    file_name="presentation.pdf",
                    mime="application/pdf",
                    use_container_width=True
                )
